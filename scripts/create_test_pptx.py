"""Create test PPTX fixtures programmatically."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

FIXTURES_DIR = Path(__file__).parent.parent / "tests" / "fixtures"


def create_bad_pptx(output_path: Path | str) -> None:
    """Create a PPTX with deliberate formatting violations."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide_layout = prs.slide_layouts[1]  # Title and Content

    # Slide 1: Title slide (no violations expected for slide number)
    slide1 = prs.slides.add_slide(slide_layout)
    title = slide1.shapes.title
    title.text = "测试演示文稿"
    # Wrong font on title
    for para in title.text_frame.paragraphs:
        for run in para.runs:
            run.font.name = "宋体"  # Should be 微软雅黑
            run.font.size = Pt(32)  # Should be 28
            run.font.color.rgb = RGBColor(0xFF, 0x00, 0x00)  # Wrong color
            run.font.bold = False  # Should be True

    # Slide 2: Content slide with violations
    slide2 = prs.slides.add_slide(slide_layout)
    title2 = slide2.shapes.title
    title2.text = "第二页标题"

    # Title with wrong font
    for para in title2.text_frame.paragraphs:
        for run in para.runs:
            run.font.name = "黑体"  # Wrong
            run.font.size = Pt(24)  # Wrong
            run.font.bold = False  # Wrong

    # Body text with wrong alignment
    body = slide2.placeholders[1]
    tf = body.text_frame
    tf.text = "这是正文内容，使用了错误的字体和对齐方式。"
    for para in tf.paragraphs:
        para.alignment = PP_ALIGN.CENTER  # Should be LEFT
        for run in para.runs:
            run.font.name = "Arial"  # Wrong font
            run.font.size = Pt(18)  # Wrong size
            run.font.color.rgb = RGBColor(0x00, 0xFF, 0x00)  # Wrong color

    # Slide 3: Missing slide number (no shapes in bottom area with numbers)
    slide3 = prs.slides.add_slide(slide_layout)
    title3 = slide3.shapes.title
    title3.text = "第三页 — 无页码"

    # Add a shape with wrong background color
    from pptx.enum.shapes import MSO_SHAPE
    shape = slide3.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(1), Inches(2), Inches(3), Inches(2),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0x00)  # Not in allowed list
    tf = shape.text_frame
    tf.text = "颜色违规的形状"

    output_path = Path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))


def create_good_pptx(output_path: Path | str) -> None:
    """Create a PPTX that should pass all checks (except slide number on slide 1)."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide_layout = prs.slide_layouts[1]  # Title and Content

    # Slide 1: Title slide
    slide1 = prs.slides.add_slide(slide_layout)
    title = slide1.shapes.title
    title.text = "规范的演示文稿"
    for para in title.text_frame.paragraphs:
        para.alignment = PP_ALIGN.LEFT
        for run in para.runs:
            run.font.name = "微软雅黑"
            run.font.size = Pt(28)
            run.font.bold = True
            run.font.color.rgb = RGBColor(0x1F, 0x2D, 0x3D)

    # Slide 2: Content slide with proper formatting
    slide2 = prs.slides.add_slide(slide_layout)
    title2 = slide2.shapes.title
    title2.text = "规范的第二页"
    for para in title2.text_frame.paragraphs:
        para.alignment = PP_ALIGN.LEFT
        for run in para.runs:
            run.font.name = "微软雅黑"
            run.font.size = Pt(28)
            run.font.bold = True
            run.font.color.rgb = RGBColor(0x1F, 0x2D, 0x3D)

    body = slide2.placeholders[1]
    tf = body.text_frame
    tf.text = "这是规范的正文内容。"
    for para in tf.paragraphs:
        para.alignment = PP_ALIGN.LEFT
        for run in para.runs:
            run.font.name = "微软雅黑"
            run.font.size = Pt(14)
            run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    # Add slide number on slide 2
    sn_shape = slide2.shapes.add_textbox(Inches(12), Inches(6.8), Inches(1), Inches(0.5))
    tf = sn_shape.text_frame
    tf.text = "2"

    output_path = Path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))


if __name__ == "__main__":
    FIXTURES_DIR.mkdir(parents=True, exist_ok=True)
    create_bad_pptx(FIXTURES_DIR / "bad.pptx")
    create_good_pptx(FIXTURES_DIR / "good.pptx")
    print(f"Created fixtures in {FIXTURES_DIR}")
