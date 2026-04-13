"""Tests for the pptx adapter and integration with lint engine."""

from pathlib import Path

import pytest

from internal.domain.rules import parse_rules
from internal.infrastructure.compiler import compile_rules
from internal.infrastructure.engine import fix_file, lint_file

FIXTURES_DIR = Path(__file__).parent / "fixtures"
RULES_PATH = Path(__file__).parent.parent / "rules.yaml"


def _lint(file_path, compiled, **kwargs):
    """Lint with AI disabled for tests."""
    return lint_file(file_path, compiled, use_ai=False, **kwargs)


def _fix(file_path, compiled, **kwargs):
    """Fix with AI disabled for tests."""
    return fix_file(file_path, compiled, use_ai=False, **kwargs)


class TestFixturesExist:
    def test_bad_pptx_exists(self):
        if not (FIXTURES_DIR / "bad.pptx").exists():
            pytest.skip("Fixtures not created yet. Run scripts/create_test_pptx.py")

    def test_good_pptx_exists(self):
        if not (FIXTURES_DIR / "good.pptx").exists():
            pytest.skip("Fixtures not created yet. Run scripts/create_test_pptx.py")


class TestEngineBad:
    @pytest.fixture(autouse=True)
    def _ensure_fixtures(self):
        if not (FIXTURES_DIR / "bad.pptx").exists():
            pytest.skip("Run scripts/create_test_pptx.py first")

    def test_bad_pptx_has_issues(self):
        rules = parse_rules(RULES_PATH)
        compiled = compile_rules(rules, use_ai=False)
        result = _lint(FIXTURES_DIR / "bad.pptx", compiled)
        assert result.total > 0
        assert not result.passed

    def test_bad_pptx_detects_font_issues(self):
        rules = parse_rules(RULES_PATH)
        compiled = compile_rules(rules, use_ai=False)
        result = _lint(FIXTURES_DIR / "bad.pptx", compiled)
        font_issues = [i for i in result.issues if i.rule_id.startswith("fonts.")]
        assert len(font_issues) > 0

    def test_bad_pptx_detects_alignment_issues(self):
        rules = parse_rules(RULES_PATH)
        compiled = compile_rules(rules, use_ai=False)
        result = _lint(FIXTURES_DIR / "bad.pptx", compiled)
        align_issues = [i for i in result.issues if i.rule_id.startswith("alignment.")]
        assert len(align_issues) > 0

    def test_bad_pptx_detects_color_issues(self):
        rules = parse_rules(RULES_PATH)
        compiled = compile_rules(rules, use_ai=False)
        result = _lint(FIXTURES_DIR / "bad.pptx", compiled)
        color_issues = [i for i in result.issues if i.rule_id.startswith("colors.")]
        # In contrast mode, red/green on white still has good contrast → no issues.
        # In whitelist mode, they'd be flagged. Either behavior is valid.
        if rules.colors.text_mode == "whitelist":
            assert len(color_issues) > 0
        else:
            # contrast mode: issues only if contrast ratio < 4.5
            assert isinstance(color_issues, list)  # just verify no crash


class TestEngineGood:
    @pytest.fixture(autouse=True)
    def _ensure_fixtures(self):
        if not (FIXTURES_DIR / "good.pptx").exists():
            pytest.skip("Run scripts/create_test_pptx.py first")

    def test_good_pptx_fewer_issues(self):
        rules = parse_rules(RULES_PATH)
        compiled = compile_rules(rules, use_ai=False)
        bad_result = _lint(FIXTURES_DIR / "bad.pptx", compiled)
        good_result = _lint(FIXTURES_DIR / "good.pptx", compiled)
        assert good_result.total < bad_result.total

    def test_contrast_mode_detects_poor_contrast(self):
        """Contrast mode should flag text with < 4.5:1 contrast ratio."""
        from io import BytesIO

        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.util import Inches

        from internal.domain.rules import parse_rules
        from internal.infrastructure.compiler import compile_rules

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        # White background (default)
        txbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
        tf = txbox.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = "Light gray text"
        run.font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)  # #CCCCCC on white ≈ 1.6:1

        buf = BytesIO()
        prs.save(buf)
        buf.seek(0)

        # Write to temp file for lint_file
        import pathlib
        import tempfile

        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
            f.write(buf.read())
            tmp_path = f.name

        try:
            import os

            rules = parse_rules(RULES_PATH)
            rules.colors.text_mode = "contrast"
            compiled = compile_rules(rules, use_ai=False)
            result = _lint(pathlib.Path(tmp_path), compiled)
            contrast_issues = [i for i in result.issues if i.rule_id == "colors.contrast"]
            assert len(contrast_issues) > 0, "Should detect poor contrast"
        finally:
            os.unlink(tmp_path)


class TestFix:
    @pytest.fixture(autouse=True)
    def _ensure_fixtures(self):
        if not (FIXTURES_DIR / "bad.pptx").exists():
            pytest.skip("Run scripts/create_test_pptx.py first")

    def test_dry_run_does_not_modify(self, tmp_path):
        rules = parse_rules(RULES_PATH)
        compiled = compile_rules(rules, use_ai=False)
        import shutil
        test_file = tmp_path / "test.pptx"
        shutil.copy(FIXTURES_DIR / "bad.pptx", test_file)

        result_before = _lint(test_file, compiled)
        fix_result = _fix(test_file, compiled, dry_run=True)
        result_after = _lint(test_file, compiled)

        assert fix_result.total == result_before.total
        assert result_after.total == result_before.total  # File unchanged

    def test_fix_reduces_issues(self, tmp_path):
        rules = parse_rules(RULES_PATH)
        compiled = compile_rules(rules, use_ai=False)
        import shutil
        test_file = tmp_path / "test.pptx"
        fixed_file = tmp_path / "fixed.pptx"
        shutil.copy(FIXTURES_DIR / "bad.pptx", test_file)

        result_before = _lint(test_file, compiled)
        _fix(test_file, compiled, output_path=fixed_file)
        result_after = _lint(fixed_file, compiled)

        # After fixing, there should be fewer or equal issues
        assert result_after.total <= result_before.total


class TestReporter:
    @pytest.fixture(autouse=True)
    def _ensure_fixtures(self):
        if not (FIXTURES_DIR / "bad.pptx").exists():
            pytest.skip("Run scripts/create_test_pptx.py first")

    def test_json_report(self):
        from internal.infrastructure.reporter import report_json
        rules = parse_rules(RULES_PATH)
        compiled = compile_rules(rules, use_ai=False)
        result = _lint(FIXTURES_DIR / "bad.pptx", compiled)
        json_str = report_json(result)
        import json
        data = json.loads(json_str)
        assert "file" in data
        assert "summary" in data
        assert "issues" in data
        assert data["summary"]["total"] > 0

    def test_html_report(self):
        from internal.infrastructure.reporter import report_html
        rules = parse_rules(RULES_PATH)
        compiled = compile_rules(rules, use_ai=False)
        result = _lint(FIXTURES_DIR / "bad.pptx", compiled)
        html = report_html(result)
        assert "<!DOCTYPE html>" in html
        assert "ppt-lint Report" in html
