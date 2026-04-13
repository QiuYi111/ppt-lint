"""Rule compiler — compiles RuleSet into check functions.

Primitive rules → direct python-pptx check/fix functions
AI rules → Claude API calls, cached as Python functions

Role classification is delegated to the engine (via role_classifier),
which passes a pre-computed role_map to each checker.
"""

from __future__ import annotations

import hashlib
import importlib.util
import logging
import os
from collections.abc import Callable
from pathlib import Path
from typing import Any

from internal.domain.models import (
    CompiledRuleSet,
    FixAction,
    FontRule,
    LintIssue,
    RuleSet,
    Severity,
)

from .pptx_adapter import (
    _emu_to_pt,
    apply_alignment_fix,
    apply_line_spacing_fix,
    get_chart_shapes,
    get_slide_background_color,
    get_slide_number_shapes,
    get_text_runs,
    hex_to_rgb,
)

logger = logging.getLogger(__name__)

CACHE_DIR = Path(".ppt-lint-cache")


def compile_rules(rules: RuleSet, use_ai: bool = True) -> CompiledRuleSet:
    """Compile a RuleSet into a CompiledRuleSet of checker functions."""
    checkers: list[Callable] = []
    user_roles_list = list(rules.fonts.keys()) if rules.fonts else []

    # Primitive rules
    checkers.extend(_compile_font_rules(rules))
    checkers.extend(_compile_color_rules(rules))
    checkers.extend(_compile_alignment_rules(rules))
    checkers.extend(_compile_spacing_rules(rules))
    checkers.extend(_compile_slide_number_rules(rules))
    checkers.extend(_compile_chart_rules(rules))

    # AI rules
    if use_ai and rules.ai_rules:
        checkers.extend(_compile_ai_rules(rules))

    return CompiledRuleSet(checkers=checkers, user_roles=user_roles_list)


# ── Font Rules ──────────────────────────────────────────

def _compile_font_rules(rules: RuleSet) -> list[Callable]:
    checkers = []
    for role, font_rule in rules.fonts.items():
        if any([font_rule.family, font_rule.size_pt, font_rule.bold is not None, font_rule.color]):
            checkers.append(_make_font_checker(role, font_rule))
    return checkers


def _make_font_checker(role: str, fr: FontRule) -> Callable:
    """Create a font checker that uses pre-computed role_map."""

    def check(
        slide: Any, slide_index: int, role_map: dict[int, str] | None = None,
    ) -> list[LintIssue]:
        issues = []
        for run_info in get_text_runs(slide):
            si = run_info["shape_index"]

            # Use AI role map if available, otherwise skip (engine handles fallback)
            text_role = None
            if role_map is not None:
                text_role = role_map.get(si)

            if text_role != role:
                continue

            msgs = []
            fixes = []

            # Check family
            if fr.family and run_info["font_name"] and run_info["font_name"] != fr.family:
                msgs.append(f"字体应为 '{fr.family}'，实际为 '{run_info['font_name']}'")
                fixes.append(("family", fr.family))

            # Check size
            if fr.size_pt is not None and run_info["font_size_pt"] is not None:
                if abs(run_info["font_size_pt"] - fr.size_pt) > 0.5:
                    msgs.append(f"字号应为 {fr.size_pt}pt，实际为 {run_info['font_size_pt']:.1f}pt")
                    fixes.append(("size_pt", fr.size_pt))

            # Check bold
            if fr.bold is not None and run_info["bold"] is not None and run_info["bold"] != fr.bold:
                expected = "加粗" if fr.bold else "不加粗"
                actual = "加粗" if run_info["bold"] else "不加粗"
                msgs.append(f"应为{expected}，实际为{actual}")
                fixes.append(("bold", fr.bold))

            # Check color
            if fr.color and run_info["color_hex"]:
                if run_info["color_hex"].upper() != fr.color.upper():
                    msgs.append(f"颜色应为 '{fr.color}'，实际为 '{run_info['color_hex']}'")
                    fixes.append(("color", fr.color))

            if msgs:
                desc = f"{role}文字 \"{run_info['text'][:20]}\""
                message = "；".join(msgs)

                fix = None
                if fixes:
                    pi = run_info["para_index"]
                    ri = run_info["run_index"]
                    fix_desc = f"修正 {desc} 的字体格式"
                    fix = FixAction(
                        fix_desc,
                        lambda s=slide, _si=si, _pi=pi, _ri=ri, _f=fixes:
                            _apply_font_fixes(s, _si, _pi, _ri, _f),
                    )

                issues.append(LintIssue(
                    rule_id=f"fonts.{role}",
                    severity=Severity.ERROR,
                    slide_index=slide_index,
                    element_desc=desc,
                    message=message,
                    fix=fix,
                ))
        return issues

    return check


def _apply_font_fixes(slide: Any, si: int, pi: int, ri: int, fixes: list[tuple]) -> None:
    shape = slide.shapes[si]
    para = shape.text_frame.paragraphs[pi]
    run = para.runs[ri]
    for key, value in fixes:
        if key == "family":
            run.font.name = value
        elif key == "size_pt":
            from pptx.util import Pt
            run.font.size = Pt(value)
        elif key == "bold":
            run.font.bold = value
        elif key == "color":
            run.font.color.rgb = hex_to_rgb(value)


# ── Color Rules ──────────────────────────────────────────

def _compile_color_rules(rules: RuleSet) -> list[Callable]:
    checkers = []
    cr = rules.colors
    if cr.allowed_text:
        checkers.append(_make_text_color_checker(cr.allowed_text, cr.text_mode))
    if cr.allowed_background:
        checkers.append(_make_bg_color_checker(cr.allowed_background))
    if cr.accent:
        pass
    return checkers


def _make_text_color_checker(allowed: list[str], mode: str = "whitelist") -> Callable:
    allowed_upper = {c.upper().lstrip("#") for c in allowed}

    def check(
        slide: Any, slide_index: int, role_map: dict[int, str] | None = None,
    ) -> list[LintIssue]:
        issues = []

        if mode == "contrast":
            # Contrast mode: check WCAG contrast ratio against slide background
            bg_color = get_slide_background_color(slide)
            # Default to white if no explicit background
            bg_rgb = _hex_to_rgb_tuple(bg_color) if bg_color else (255, 255, 255)

            for run_info in get_text_runs(slide):
                color_hex = run_info["color_hex"]
                if not color_hex:
                    continue
                text_rgb = _hex_to_rgb_tuple(color_hex)
                if text_rgb is None:
                    continue
                ratio = _contrast_ratio(bg_rgb, text_rgb)
                if ratio < 4.5:
                    desc = f"文字 \"{run_info['text'][:20]}\""
                    issues.append(LintIssue(
                        rule_id="colors.contrast",
                        severity=Severity.WARNING,
                        slide_index=slide_index,
                        element_desc=desc,
                        message=f"文字与背景对比度不足 ({ratio:.1f}:1，建议 ≥ 4.5:1)",
                        fix=None,
                    ))
        else:
            # Whitelist mode: exact color match
            for run_info in get_text_runs(slide):
                color_hex = run_info["color_hex"]
                if not color_hex:
                    continue
                color_clean = color_hex.upper().lstrip("#")
                if color_clean not in allowed_upper:
                    desc = f"文字 \"{run_info['text'][:20]}\""
                    issues.append(LintIssue(
                        rule_id="colors.allowed_text",
                        severity=Severity.WARNING,
                        slide_index=slide_index,
                        element_desc=desc,
                        message=f"文字颜色 '{color_hex}' 不在允许列表中",
                        fix=None,
                    ))
        return issues

    return check


def _hex_to_rgb_tuple(hex_str: str) -> tuple[int, int, int] | None:
    """Convert '#1F2D3D' to (31, 45, 61)."""
    try:
        h = hex_str.lstrip("#")
        return (int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))
    except (ValueError, IndexError):
        return None


def _luminance(r: int, g: int, b: int) -> float:
    """Calculate relative luminance per WCAG 2.0."""
    def _linearize(c: int) -> float:
        s = c / 255.0
        return s / 12.92 if s <= 0.03928 else ((s + 0.055) / 1.055) ** 2.4
    return 0.2126 * _linearize(r) + 0.7152 * _linearize(g) + 0.0722 * _linearize(b)


def _contrast_ratio(
    rgb1: tuple[int, int, int], rgb2: tuple[int, int, int],
) -> float:
    """Calculate WCAG contrast ratio between two RGB colors."""
    l1 = _luminance(*rgb1)
    l2 = _luminance(*rgb2)
    lighter = max(l1, l2)
    darker = min(l1, l2)
    return (lighter + 0.05) / (darker + 0.05)


def _make_bg_color_checker(allowed: list[str]) -> Callable:
    allowed_upper = {c.upper().lstrip("#") for c in allowed}

    def check(
        slide: Any, slide_index: int, role_map: dict[int, str] | None = None,
    ) -> list[LintIssue]:
        issues = []
        bg_color = get_slide_background_color(slide)
        if bg_color:
            color_clean = bg_color.upper().lstrip("#")
            if color_clean not in allowed_upper:
                issues.append(LintIssue(
                    rule_id="colors.allowed_background",
                    severity=Severity.WARNING,
                    slide_index=slide_index,
                    element_desc="幻灯片背景",
                    message=f"背景色 '{bg_color}' 不在允许列表中",
                    fix=FixAction(
                        "修正幻灯片背景色为第一个允许色",
                        lambda s=slide, c=allowed[0]: _set_slide_bg_color(s, c),
                    ),
                ))
        return issues

    return check


def _set_slide_bg_color(slide: Any, hex_color: str) -> None:
    """Set slide background color."""
    try:
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = hex_to_rgb(hex_color)
    except (AttributeError, TypeError):
        pass


# ── Alignment Rules ──────────────────────────────────────

def _compile_alignment_rules(rules: RuleSet) -> list[Callable]:
    checkers = []
    ar = rules.alignment
    if ar.title:
        checkers.append(_make_alignment_checker("title", ar.title))
    if ar.body:
        checkers.append(_make_alignment_checker("body", ar.body))
    return checkers


def _make_alignment_checker(role: str, expected_align: str) -> Callable:

    def check(
        slide: Any, slide_index: int, role_map: dict[int, str] | None = None,
    ) -> list[LintIssue]:
        issues = []
        for run_info in get_text_runs(slide):
            si = run_info["shape_index"]

            text_role = None
            if role_map is not None:
                text_role = role_map.get(si)

            if text_role != role:
                continue

            actual = run_info["paragraph_alignment"]
            if actual is None:
                continue

            if actual != expected_align:
                desc = f"{role}段落"
                pi = run_info["para_index"]
                fix = FixAction(
                    f"修正 {desc} 对齐为 {expected_align}",
                    lambda s=slide, _si=si, _pi=pi, a=expected_align:
                        apply_alignment_fix(s.shapes[_si].text_frame.paragraphs[_pi], a),
                )
                issues.append(LintIssue(
                    rule_id=f"alignment.{role}",
                    severity=Severity.WARNING,
                    slide_index=slide_index,
                    element_desc=desc,
                    message=f"对齐应为 '{expected_align}'，实际为 '{actual}'",
                    fix=fix,
                ))
        return issues

    return check


# ── Spacing Rules ────────────────────────────────────────

def _compile_spacing_rules(rules: RuleSet) -> list[Callable]:
    checkers = []
    sr = rules.spacing
    if sr.line_spacing is not None:
        checkers.append(_make_line_spacing_checker(sr.line_spacing))
    return checkers


def _make_line_spacing_checker(expected: float) -> Callable:

    def check(
        slide: Any, slide_index: int, role_map: dict[int, str] | None = None,
    ) -> list[LintIssue]:
        issues = []
        for run_info in get_text_runs(slide):
            ls = run_info["line_spacing"]
            if ls is None:
                continue
            actual_pt = _emu_to_pt(ls)
            if actual_pt is None:
                continue
            expected_pt = expected * 12
            if abs(actual_pt - expected_pt) > 1.0:
                desc = f"段落 \"{run_info['text'][:20]}\""
                si, pi = run_info["shape_index"], run_info["para_index"]
                fix = FixAction(
                    f"修正 {desc} 行间距",
                    lambda s=slide, _si=si, _pi=pi, e=expected:
                        apply_line_spacing_fix(
                            s.shapes[_si].text_frame.paragraphs[_pi], e,
                        ),
                )
                issues.append(LintIssue(
                    rule_id="spacing.line_spacing",
                    severity=Severity.WARNING,
                    slide_index=slide_index,
                    element_desc=desc,
                    message=f"行间距偏差较大 (实际 ~{actual_pt:.1f}pt)",
                    fix=fix,
                ))
        return issues

    return check


# ── Slide Number Rules ───────────────────────────────────

def _compile_slide_number_rules(rules: RuleSet) -> list[Callable]:
    checkers = []
    snr = rules.slide_number
    if snr.visible:
        checkers.append(_make_slide_number_visible_checker())
    return checkers


def _make_slide_number_visible_checker() -> Callable:
    def check(
        slide: Any, slide_index: int, role_map: dict[int, str] | None = None,
    ) -> list[LintIssue]:
        if slide_index == 0:
            return []

        # First try role_map — if any shape is classified as slide_number
        if role_map is not None:
            has_sn = any(r == "slide_number" for r in role_map.values())
            if has_sn:
                return []

        # Fallback: heuristic detection
        sn_shapes = get_slide_number_shapes(slide)
        if not sn_shapes:
            return [LintIssue(
                rule_id="slide_number.visible",
                severity=Severity.WARNING,
                slide_index=slide_index,
                element_desc="页码",
                message="未找到页码",
                fix=None,
            )]
        return []
    return check


# ── Chart Rules ──────────────────────────────────────────

def _compile_chart_rules(rules: RuleSet) -> list[Callable]:
    checkers = []
    cr = rules.charts
    if cr.require_title:
        checkers.append(_make_chart_title_checker())
    if cr.require_axis_labels:
        checkers.append(_make_chart_axis_checker())
    return checkers


def _make_chart_title_checker() -> Callable:
    def check(
        slide: Any, slide_index: int, role_map: dict[int, str] | None = None,
    ) -> list[LintIssue]:
        issues = []
        for chart_info in get_chart_shapes(slide):
            if not chart_info["has_title"]:
                desc = f"图表 \"{chart_info['shape_name']}\""
                issues.append(LintIssue(
                    rule_id="charts.require_title",
                    severity=Severity.WARNING,
                    slide_index=slide_index,
                    element_desc=desc,
                    message="图表缺少标题",
                    fix=None,
                ))
        return issues
    return check


def _make_chart_axis_checker() -> Callable:
    def check(
        slide: Any, slide_index: int, role_map: dict[int, str] | None = None,
    ) -> list[LintIssue]:
        issues = []
        for chart_info in get_chart_shapes(slide):
            desc = f"图表 \"{chart_info['shape_name']}\""
            try:
                chart = slide.shapes[chart_info["shape_index"]].chart
                cat_axis = chart.category_axis
                if cat_axis and not cat_axis.has_title:
                    issues.append(LintIssue(
                        rule_id="charts.require_axis_labels",
                        severity=Severity.INFO,
                        slide_index=slide_index,
                        element_desc=desc,
                        message="图表 X 轴缺少标签",
                        fix=None,
                    ))
                val_axis = chart.value_axis
                if val_axis and not val_axis.has_title:
                    issues.append(LintIssue(
                        rule_id="charts.require_axis_labels",
                        severity=Severity.INFO,
                        slide_index=slide_index,
                        element_desc=desc,
                        message="图表 Y 轴缺少标签",
                        fix=None,
                    ))
            except Exception:
                pass
        return issues
    return check


# ── AI Rules ─────────────────────────────────────────────

def _compile_ai_rules(rules: RuleSet) -> list[Callable]:
    checkers = []
    for ai_rule in rules.ai_rules:
        checker = _load_or_compile_ai_rule(ai_rule)
        if checker:
            checkers.append(checker)
    return checkers


def _load_or_compile_ai_rule(ai_rule: Any) -> Callable | None:
    """Load cached AI rule or compile via Claude API."""
    CACHE_DIR.mkdir(exist_ok=True)

    content = f"{ai_rule.id}:{ai_rule.description}:{ai_rule.severity.value}"
    hash_key = hashlib.sha256(content.encode()).hexdigest()[:12]
    cache_file = CACHE_DIR / f"{ai_rule.id}_{hash_key}.py"
    hash_file = CACHE_DIR / f"{ai_rule.id}.hash"

    if cache_file.exists():
        stored_hash = hash_file.read_text().strip() if hash_file.exists() else ""
        if stored_hash == hash_key:
            try:
                return _load_cached_function(cache_file, ai_rule.id)
            except Exception as e:
                logger.warning(f"Failed to load cached rule {ai_rule.id}: {e}")

    fn = _claude_compile_rule(ai_rule)
    if fn:
        cache_file.write_text(_generate_cache_module(ai_rule), encoding="utf-8")
        hash_file.write_text(hash_key, encoding="utf-8")
        return fn

    return None


def _claude_compile_rule(ai_rule: Any) -> Callable | None:
    """Call Claude API to compile an AI rule into a Python function."""
    try:
        import anthropic
    except ImportError:
        logger.warning("anthropic package not installed, skipping AI rules")
        return None

    api_key = os.environ.get("ANTHROPIC_API_KEY")
    if not api_key:
        logger.warning("ANTHROPIC_API_KEY not set, skipping AI rules")
        return None

    prompt = f"""你是一个 python-pptx 专家。请将以下规则描述编译为一个 Python 函数。
规则描述：{ai_rule.description}
严重级别：{ai_rule.severity.value}
函数签名：
  def check_{ai_rule.id}(slide, slide_index) -> list:
      ...
返回的列表中每个元素是一个字典，包含以下键：
  rule_id: str (使用 "{ai_rule.id}")
  severity: str ("{ai_rule.severity.value}")
  slide_index: int
  element_desc: str
  message: str
  fix: None

slide 参数是 python-pptx 的 Slide 对象。
只返回 Python 代码，不要解释。不要使用 import 语句（pptx 模块已经可用）。"""

    try:
        client = anthropic.Anthropic(api_key=api_key)
        response = client.messages.create(
            model="claude-sonnet-4-20250514",
            max_tokens=2000,
            messages=[{"role": "user", "content": prompt}],
        )
        code = response.content[0].text.strip()
        if "```python" in code:
            code = code.split("```python")[1].split("```")[0].strip()
        elif "```" in code:
            code = code.split("```")[1].split("```")[0].strip()

        ns: dict[str, Any] = {}
        exec(code, ns)  # noqa: S102
        fn = ns.get(f"check_{ai_rule.id}")
        if fn and callable(fn):
            return fn
        return None
    except Exception as e:
        logger.error(f"Failed to compile AI rule {ai_rule.id}: {e}")
        return None


def _generate_cache_module(ai_rule: Any) -> str:
    """Generate the cache .py file content for an AI rule."""
    return f"""# Auto-generated by ppt-lint AI compiler
# Rule: {ai_rule.id}
# Description: {ai_rule.description}
# Severity: {ai_rule.severity.value}
#
# This file is managed by ppt-lint. Do not edit manually.
# To regenerate, delete the .ppt-lint-cache/ directory.

from internal.domain.models import LintIssue, Severity

# The compiled function will be loaded dynamically.
# This is a placeholder — the actual function is stored in the compiled cache.
"""


def _load_cached_function(cache_file: Path, rule_id: str) -> Callable | None:
    """Load a compiled function from cache."""
    spec = importlib.util.spec_from_file_location(f"cache_{rule_id}", str(cache_file))
    if spec and spec.loader:
        mod = importlib.util.module_from_spec(spec)
        spec.loader.exec_module(mod)
        fn = getattr(mod, f"check_{rule_id}", None)
        if callable(fn):
            return fn
    return None
