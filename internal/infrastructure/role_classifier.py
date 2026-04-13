"""AI-powered text role classifier via Anthropic-compatible API.

Instead of brittle heuristics, sends slide shape metadata to an LLM
for intelligent role classification. Uses direct HTTP calls to an
Anthropic-compatible API (e.g. Z.AI proxy) instead of Claude CLI,
which has compatibility issues with some proxy backends.

Supports two modes:
  - Batch (default): all slides classified in one API call
  - Per-slide fallback: one call per slide (for cache hits / retries)

Results are cached per slide content hash.
"""

from __future__ import annotations

import hashlib
import json
import logging
import os
import time
import urllib.error
import urllib.request
from pathlib import Path
from typing import Any

from pptx import Presentation

from .pptx_adapter import classify_text_role, extract_slide_summary

logger = logging.getLogger(__name__)

CACHE_DIR = Path(".ppt-lint-cache/roles")

# API config — reads from env vars with sensible defaults
_API_BASE_URL = os.environ.get(
    "PPT_LINT_API_BASE_URL",
    os.environ.get("ANTHROPIC_BASE_URL", "https://api.z.ai/api/anthropic"),
)
_API_KEY = os.environ.get(
    "PPT_LINT_API_KEY",
    os.environ.get("ANTHROPIC_AUTH_TOKEN", ""),
)
_API_MODEL = os.environ.get(
    "PPT_LINT_MODEL",
    os.environ.get("ANTHROPIC_DEFAULT_OPUS_MODEL", "glm-5.1"),
)
_API_TIMEOUT = int(os.environ.get("PPT_LINT_API_TIMEOUT", "120"))

# Proxy settings (inherit from env)
_HTTP_PROXY = os.environ.get("HTTPS_PROXY") or os.environ.get("HTTP_PROXY", "")


def _get_opener() -> urllib.request.OpenerDirector:
    """Build URL opener, optionally with proxy."""
    if _HTTP_PROXY:
        proxy = urllib.request.ProxyHandler({
            "https": _HTTP_PROXY,
            "http": _HTTP_PROXY,
        })
        return urllib.request.build_opener(proxy)
    return urllib.request.build_opener()


def _call_api(prompt: str, max_tokens: int = 4096) -> str | None:
    """Send a prompt to the Anthropic-compatible API and return the text response.

    Returns None on any error (timeout, HTTP error, parse error).
    Retries once on 429 (rate limit) with a 5s backoff.
    """
    url = f"{_API_BASE_URL.rstrip('/')}/v1/messages"
    payload = json.dumps({
        "model": _API_MODEL,
        "max_tokens": max_tokens,
        "messages": [{"role": "user", "content": prompt}],
    }).encode()

    req = urllib.request.Request(url, data=payload, headers={
        "Content-Type": "application/json",
        "x-api-key": _API_KEY,
        "anthropic-version": "2023-06-01",
    })

    opener = _get_opener()

    for attempt in range(2):
        try:
            resp = opener.open(req, timeout=_API_TIMEOUT)
            data = json.loads(resp.read().decode())
            # Extract text from Anthropic response format
            blocks = data.get("content", [])
            for block in blocks:
                if block.get("type") == "text":
                    return block["text"]
            return None

        except urllib.error.HTTPError as e:
            if e.code == 429 and attempt == 0:
                retry_after = e.headers.get("retry-after")
                wait = int(retry_after) if retry_after else 5
                logger.warning(f"API rate limited, retrying in {wait}s...")
                time.sleep(wait)
                continue
            body = e.read().decode()[:300] if hasattr(e, "read") else ""
            logger.warning(f"API HTTP {e.code}: {body}")
            return None

        except urllib.error.URLError as e:
            logger.warning(f"API connection error: {e.reason}")
            return None
        except TimeoutError:
            logger.warning(f"API timed out ({_API_TIMEOUT}s)")
            return None
        except (json.JSONDecodeError, KeyError) as e:
            logger.warning(f"API response parse error: {e}")
            return None
        except Exception as e:
            logger.warning(f"API call error: {e}")
            return None

    return None


def _api_available() -> bool:
    """Check if API key is configured."""
    return bool(_API_KEY)


def _parse_json_response(content: str) -> dict[str, Any] | None:
    """Parse JSON from LLM response, stripping markdown fences if present."""
    if not content:
        return None

    # Strip markdown fences
    if "```" in content:
        content = content.split("```")[1]
        if content.startswith("json"):
            content = content[4:]
        content = content.strip()

    try:
        return json.loads(content)
    except json.JSONDecodeError:
        return None


# ─── Public API ─────────────────────────────────────────────────────────


def classify_all_slides(
    prs: Presentation,
    user_roles: set[str] | None = None,
    use_ai: bool = True,
) -> dict[int, dict[int, str]]:
    """Classify text shapes on all slides.

    Returns dict mapping slide_index → {shape_index → role_name}.
    """
    n_slides = len(prs.slides)
    result: dict[int, dict[int, str]] = {}

    if not use_ai or not _api_available():
        logger.info("AI disabled or no API key — using heuristic classifier")
        for i, slide in enumerate(prs.slides):
            result[i] = _classify_heuristic(slide, user_roles)
        return result

    # Try batch mode first: one API call for all slides
    batch_result = _classify_batch(prs, user_roles)
    if batch_result is not None:
        for i in range(n_slides):
            if i in batch_result:
                result[i] = batch_result[i]
            else:
                result[i] = _classify_heuristic(prs.slides[i], user_roles)
        return result

    # Fallback: per-slide calls
    logger.warning("Batch classification failed, falling back to per-slide")
    for i, slide in enumerate(prs.slides):
        result[i] = classify_slide_roles(slide, i, user_roles, use_ai=True)
    return result


def classify_slide_roles(
    slide: Any,
    slide_index: int,
    user_roles: set[str] | None = None,
    use_ai: bool = True,
) -> dict[int, str]:
    """Classify all text shapes on a single slide into roles.

    Returns a dict mapping shape_index → role_name.
    """
    if not use_ai:
        return _classify_heuristic(slide, user_roles)

    # Check per-slide cache first
    summary = extract_slide_summary(slide, slide_index)
    text_shapes = [s for s in summary["shapes"] if s.get("has_text")]
    if not text_shapes:
        return {}

    cache_key = _slide_content_hash(text_shapes, user_roles)
    cached = _load_cache(cache_key)
    if cached is not None:
        logger.debug(f"Role cache hit for slide {slide_index}")
        return cached

    result = _classify_slide_via_api(slide, slide_index, user_roles)
    if result is not None:
        return result

    logger.info(
        "API unavailable, falling back to heuristic "
        f"classifier for slide {slide_index}"
    )
    return _classify_heuristic(slide, user_roles)


# ─── Heuristic classifier ───────────────────────────────────────────────


def _classify_heuristic(
    slide: Any, user_roles: set[str] | None = None,
) -> dict[int, str]:
    """Use the built-in heuristic classifier for all shapes."""
    roles: dict[int, str] = {}
    for si, shape in enumerate(slide.shapes):
        if shape.has_text_frame and shape.text_frame.text.strip():
            roles[si] = classify_text_role(shape, slide, user_roles)
    return roles


# ─── Batch classification ───────────────────────────────────────────────


def _classify_batch(
    prs: Presentation,
    user_roles: set[str] | None = None,
) -> dict[int, dict[int, str]] | None:
    """Classify all slides in one API call. Much faster than per-slide."""
    roles_list = sorted(user_roles) if user_roles else [
        "title", "subtitle", "body", "caption",
        "section_number", "section_title", "slide_number", "footer",
    ]

    # Build batch payload
    slides_data: list[dict[str, Any]] = []
    slides_with_text = []

    for i, slide in enumerate(prs.slides):
        summary = extract_slide_summary(slide, i)
        text_shapes = [s for s in summary["shapes"] if s.get("has_text")]
        if not text_shapes:
            continue

        slides_with_text.append(i)

        # Check per-slide cache — skip cached slides in batch
        cache_key = _slide_content_hash(text_shapes, user_roles)
        cached = _load_cache(cache_key)
        if cached is not None:
            continue  # will be filled from cache later

        # Compact shape representation for the prompt
        compact = []
        for s in text_shapes:
            entry: dict[str, Any] = {
                "i": s["index"],
                "n": s["name"],
            }
            if s.get("text_preview"):
                entry["t"] = s["text_preview"][:80]
            if s.get("max_font_size_pt") is not None:
                entry["fs"] = s["max_font_size_pt"]
            if s.get("left_in") is not None:
                entry["x"] = s["left_in"]
                entry["y"] = s["top_in"]
            if s.get("placeholder_idx") is not None:
                entry["ph"] = s["placeholder_idx"]
            compact.append(entry)

        slides_data.append({
            "slide": i,
            "w": summary["slide_width_in"],
            "h": summary["slide_height_in"],
            "shapes": compact,
        })

    if not slides_data:
        # All slides cached
        result: dict[int, dict[int, str]] = {}
        for i in slides_with_text:
            summary = extract_slide_summary(prs.slides[i], i)
            text_shapes = [s for s in summary["shapes"] if s.get("has_text")]
            cache_key = _slide_content_hash(text_shapes, user_roles)
            cached = _load_cache(cache_key)
            if cached is not None:
                result[i] = cached
        return result

    # Build prompt
    payload = json.dumps(slides_data, ensure_ascii=False, indent=1)
    n_shapes = sum(len(s["shapes"]) for s in slides_data)

    prompt = (
        "You are a PPT layout expert. Classify every text shape into ONE role.\n\n"
        f"Roles: {roles_list}\n\n"
        f"Input: {len(slides_data)} slides, {n_shapes} text shapes total.\n"
        "Each slide has: slide number, dimensions (w/h inches), shapes array.\n"
        "Each shape has: i=index, n=name, t=text preview, fs=max font size pt, "
        "x/y=position inches, ph=placeholder index.\n\n"
        "Rules:\n"
        "- Large decorative numbers (>=40pt, short text) on divider slides → section_number\n"
        "- Page numbers (e.g. '4 / 13', '2 / 25') → slide_number\n"
        "- Headers/footers (institution, date, small text at edges) → footer\n"
        "- Main heading of a content slide → title\n"
        "- Secondary/sub heading → section_title\n"
        "- Content text, bullet points, TOC entries → body\n"
        "- Figure captions, small annotations near images → caption\n"
        "- Author info, affiliation → footer\n"
        "- Unclear → body\n\n"
        f"Data:\n{payload}\n\n"
        'Return ONLY a JSON object: {"slide_index": {"shape_index": "role", ...}, ...}\n'
        "No explanation. No markdown fences."
    )

    content = _call_api(prompt, max_tokens=4096)
    if content is None:
        return None

    batch_map = _parse_json_response(content)
    if batch_map is None:
        logger.warning("Failed to parse batch classification response")
        return None

    # Parse and cache per-slide results
    all_result: dict[int, dict[int, str]] = {}

    # Fill cached slides first
    for i in slides_with_text:
        summary = extract_slide_summary(prs.slides[i], i)
        text_shapes = [s for s in summary["shapes"] if s.get("has_text")]
        cache_key = _slide_content_hash(text_shapes, user_roles)
        cached = _load_cache(cache_key)
        if cached is not None:
            all_result[i] = cached

    # Fill batch results
    roles_set = set(roles_list)
    for slide_key, shape_roles in batch_map.items():
        try:
            si = int(slide_key)
        except (ValueError, TypeError):
            continue
        parsed: dict[int, str] = {}
        for shape_key, role in shape_roles.items():
            try:
                idx = int(shape_key)
            except (ValueError, TypeError):
                continue
            if role in roles_set:
                parsed[idx] = role
            else:
                parsed[idx] = "body"
        all_result[si] = parsed
        # Cache this slide
        summary = extract_slide_summary(prs.slides[si], si)
        text_shapes = [s for s in summary["shapes"] if s.get("has_text")]
        cache_key = _slide_content_hash(text_shapes, user_roles)
        _save_cache(cache_key, parsed)

    classified_slides = len(all_result)
    classified_shapes = sum(len(v) for v in all_result.values())
    logger.info(
        f"API batch: classified {classified_shapes} shapes "
        f"across {classified_slides} slides"
    )
    return all_result


# ─── Per-slide classification ───────────────────────────────────────────


def _classify_slide_via_api(
    slide: Any,
    slide_index: int,
    user_roles: set[str] | None = None,
) -> dict[int, str] | None:
    """Classify shapes on a single slide via API (per-slide fallback)."""
    if not _api_available():
        return None

    summary = extract_slide_summary(slide, slide_index)
    text_shapes = [s for s in summary["shapes"] if s.get("has_text")]
    if not text_shapes:
        return {}

    cache_key = _slide_content_hash(text_shapes, user_roles)
    cached = _load_cache(cache_key)
    if cached is not None:
        logger.debug(f"Role cache hit for slide {slide_index}")
        return cached

    roles_list = sorted(user_roles) if user_roles else [
        "title", "subtitle", "body", "caption",
        "section_number", "section_title", "slide_number", "footer",
    ]

    # Compact representation
    compact = []
    for s in text_shapes:
        entry: dict[str, Any] = {"i": s["index"], "n": s["name"]}
        if s.get("text_preview"):
            entry["t"] = s["text_preview"][:80]
        if s.get("max_font_size_pt") is not None:
            entry["fs"] = s["max_font_size_pt"]
        if s.get("left_in") is not None:
            entry["x"] = s["left_in"]
            entry["y"] = s["top_in"]
        if s.get("placeholder_idx") is not None:
            entry["ph"] = s["placeholder_idx"]
        compact.append(entry)

    shapes_desc = json.dumps(compact, ensure_ascii=False, indent=1)

    prompt = (
        "Classify each text shape into ONE role.\n"
        f"Roles: {roles_list}\n"
        f"Slide: {summary['slide_width_in']}\"×{summary['slide_height_in']}\"\n"
        f"Shapes:\n{shapes_desc}\n\n"
        "- Large decorative numbers (>=40pt) → section_number\n"
        "- Page numbers ('4 / 13') → slide_number\n"
        "- Headers/footers → footer\n"
        "- Main heading → title\n"
        "- Sub-heading → section_title\n"
        "- Content → body\n"
        "- Captions → caption\n"
        "- Unclear → body\n\n"
        'Return JSON only: {"0": "title", "1": "body"}\n'
        "No explanation, no fences."
    )

    content = _call_api(prompt, max_tokens=1024)
    if content is None:
        return None

    role_map = _parse_json_response(content)
    if role_map is None:
        logger.warning(f"Failed to parse API response for slide {slide_index}")
        return None

    parsed: dict[int, str] = {}
    roles_set = set(roles_list)
    for k, v in role_map.items():
        try:
            idx = int(k)
            parsed[idx] = v if v in roles_set else "body"
        except (ValueError, TypeError):
            continue

    _save_cache(cache_key, parsed)
    logger.info(f"API classified {len(parsed)} shapes on slide {slide_index}")
    return parsed


# ─── Cache ──────────────────────────────────────────────────────────────


def _slide_content_hash(
    text_shapes: list[dict[str, Any]],
    user_roles: set[str] | None = None,
) -> str:
    """Create a content-based cache key."""
    relevant = []
    for s in text_shapes:
        relevant.append({
            "i": s["index"],
            "n": s["name"],
            "t": s.get("text_preview", "")[:60],
            "fs": s.get("max_font_size_pt"),
            "pos": (s.get("left_in"), s.get("top_in")),
            "ph": s.get("placeholder_idx"),
        })
    blob = json.dumps(relevant, sort_keys=True, ensure_ascii=False)
    roles_blob = json.dumps(sorted(user_roles)) if user_roles else ""
    combined = f"{blob}|{roles_blob}"
    return hashlib.sha256(combined.encode()).hexdigest()[:16]


def _load_cache(key: str) -> dict[int, str] | None:
    """Load cached role classification."""
    cache_file = CACHE_DIR / f"{key}.json"
    if cache_file.exists():
        try:
            data = json.loads(cache_file.read_text())
            return {int(k): v for k, v in data.items()}
        except (json.JSONDecodeError, OSError):
            return None
    return None


def _save_cache(key: str, roles: dict[int, str]) -> None:
    """Save role classification to cache."""
    CACHE_DIR.mkdir(parents=True, exist_ok=True)
    cache_file = CACHE_DIR / f"{key}.json"
    try:
        cache_file.write_text(
            json.dumps({str(k): v for k, v in roles.items()}, ensure_ascii=False),
            encoding="utf-8",
        )
    except OSError as e:
        logger.warning(f"Failed to write role cache: {e}")
