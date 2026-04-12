"""python-pptx adapter — wraps PPTX operations for linting and fixing."""

from __future__ import annotations

import logging
import re
from typing import Any

from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Pt

logger = logging.getLogger(__name__)

# Map alignment string names to PP_ALIGN enum
ALIGNMENT_MAP: dict[str, PP_ALIGN] = {
    "left": PP_ALIGN.LEFT,
    "center": PP_ALIGN.CENTER,
    "right": PP_ALIGN.RIGHT,
    "justify": PP_ALIGN.JUSTIFY,
}


def hex_to_rgb(hex_str: str) -> RGBColor:
    """Convert a hex color string like '#1F2D3D' to RGBColor."""
    h = hex_str.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def rgb_to_hex(rgb: RGBColor) -> str:
    """Convert RGBColor to hex string."""
    return f"#{rgb}"


def rgb_to_hex_safe(color_obj: Any) -> str | None:
    """Safely extract hex string from a color object."""
    try:
        if color_obj is None:
            return None
        rgb = color_obj.rgb
        return f"#{rgb}"
    except (AttributeError, TypeError):
        return None


def get_text_runs(slide: Any) -> list[dict[str, Any]]:
    """Extract all text runs with their formatting from a slide.

    Returns list of dicts with keys:
        shape_index, shape_name, text_frame_index, para_index, run_index,
        text, font_name, font_size_pt, bold, color_hex
    """
    runs = []
    for si, shape in enumerate(slide.shapes):
        if not shape.has_text_frame:
            continue
        for tfi, text_frame in enumerate(  # noqa: F841
            shape.text_frame.paragraphs,
        ):
            # Re-iterate properly
            pass

    # Redo properly
    for si, shape in enumerate(slide.shapes):
        if not shape.has_text_frame:
            continue
        text_frame = shape.text_frame
        for pi, para in enumerate(text_frame.paragraphs):
            for ri, run in enumerate(para.runs):
                font = run.font
                color_hex = rgb_to_hex_safe(font.color)
                runs.append({
                    "shape_index": si,
                    "shape_name": shape.name,
                    "text_frame_index": 0,
                    "para_index": pi,
                    "run_index": ri,
                    "text": run.text,
                    "font_name": font.name,
                    "font_size_pt": _emu_to_pt(font.size) if font.size else None,
                    "bold": font.bold,
                    "color_hex": color_hex,
                    "paragraph_alignment": _align_to_str(para.alignment),
                    "shape_left": shape.left,
                    "shape_top": shape.top,
                    "shape_width": shape.width,
                    "shape_height": shape.height,
                    "line_spacing": para.line_spacing,
                })
    return runs


def get_slide_background_color(slide: Any) -> str | None:
    """Get the slide's own background color (not shape fills).

    Returns the hex color of the slide background, or None if
    the slide uses a transparent/default background.
    """
    try:
        bg = slide.background
        fill = bg.fill
        if fill.type is not None:
            return rgb_to_hex_safe(fill.fore_color)
    except (AttributeError, TypeError):
        pass
    return None


def get_shapes_with_fill(slide: Any) -> list[dict[str, Any]]:
    """Get all shapes with background fill colors.

    DEPRECATED: Use get_slide_background_color() for slide background checks.
    This function returns decorative shape fills which are typically NOT
    the slide background color. Kept for backward compatibility but
    callers should prefer get_slide_background_color().
    """
    shapes = []
    for si, shape in enumerate(slide.shapes):
        fill_color = None
        try:
            if shape.fill.type is not None:
                fill_color = rgb_to_hex_safe(shape.fill.fore_color)
        except Exception:
            pass
        if fill_color:
            shapes.append({
                "shape_index": si,
                "shape_name": shape.name,
                "fill_color": fill_color,
                "shape_type": str(shape.shape_type),
            })
    return shapes


def get_slide_number_shapes(slide: Any) -> list[dict[str, Any]]:
    """Detect shapes that likely represent slide numbers."""
    candidates = []
    # Get slide width for position heuristics
    try:
        _sw = (
            slide.slide_layout.slide_master.slide_width
            if slide.slide_layout and hasattr(slide.slide_layout.slide_master, "slide_width")
            else Emu(9144000)
        )
    except AttributeError:
        _sw = Emu(9144000)
    for si, shape in enumerate(slide.shapes):
        if not shape.has_text_frame:
            continue
        text = shape.text_frame.text.strip()
        # Slide number detection: pure digit, slide_id match, or "N / M" format
        is_slide_num = False
        if text.isdigit() or text == str(slide.slide_id):
            is_slide_num = True
        elif re.match(r"^\d+\s*/\s*\d+$", text):
            is_slide_num = True

        if is_slide_num:
            candidates.append({
                "shape_index": si,
                "shape_name": shape.name,
                "text": text,
                "left": shape.left,
                "top": shape.top,
                "width": shape.width,
                "height": shape.height,
                "font_size_pt": _emu_to_pt(shape.text_frame.paragraphs[0].runs[0].font.size)
                if shape.text_frame.paragraphs and shape.text_frame.paragraphs[0].runs
                else None,
            })
    return candidates


def get_chart_shapes(slide: Any) -> list[dict[str, Any]]:
    """Get all chart shapes on a slide."""
    charts = []
    for si, shape in enumerate(slide.shapes):
        if shape.has_chart:
            chart = shape.chart
            has_title = (
                chart.has_title
                and chart.chart_title.has_text_frame
                and chart.chart_title.text_frame.text.strip()
            )
            charts.append({
                "shape_index": si,
                "shape_name": shape.name,
                "has_title": bool(has_title),
                "title_text": chart.chart_title.text_frame.text if has_title else None,
                "chart_type": str(chart.chart_type) if hasattr(chart, "chart_type") else "unknown",
            })
    return charts


def classify_text_role(shape: Any, slide: Any, user_roles: set[str] | None = None) -> str:
    """Classify a text shape into a role.

    Uses shape name, placeholder type, position, font size, and content
    to determine the role. Supports user-defined roles from rules.yaml.

    Built-in roles (matched by heuristic priority):
      - title, subtitle, section_number, slide_number, footer, caption, body

    If *user_roles* is provided, the classifier will attempt to map shapes
    to those role names; shapes that don't match any user role fall back to
    the built-in classification, ultimately defaulting to "body".

    Parameters
    ----------
    shape : pptx.shape
    slide : pptx.slide
    user_roles : set of role names defined in rules.yaml fonts section
    """
    name = shape.name.lower()
    text = shape.text_frame.text.strip() if shape.has_text_frame else ""

    # ── 1. Shape name keywords ──────────────────────────
    if "slide number" in name or "sldnum" in name or "page number" in name:
        return "slide_number"
    if "footer" in name:
        return "footer"
    if "header" in name:
        return "footer"  # headers in PPT are typically non-content text
    if "title" in name and "subtitle" not in name and "section" not in name:
        return "title"
    if "subtitle" in name:
        if user_roles and "subtitle" in user_roles:
            return "subtitle"
        return "body"
    if "section" in name or "divider" in name or "separator" in name:
        return "section_number"

    # ── 2. Placeholder index ────────────────────────────
    if hasattr(shape, "is_placeholder") and shape.is_placeholder:
        ph_fmt = getattr(shape, "placeholder_format", None)
        if ph_fmt and hasattr(ph_fmt, "idx"):
            idx = ph_fmt.idx
            if idx == 0:
                return "title"
            if idx == 1:
                # idx 1 is typically the body/subtitle placeholder.
                # Return "subtitle" if the user defined that role, else "body".
                if user_roles and "subtitle" in user_roles:
                    return "subtitle"
                return "body"
            if idx >= 10:
                # Standard PPTX: idx 10+ = date, slide number, footer, etc.
                if text.isdigit() or re.match(r"^\d+\s*/\s*\d+$", text):
                    return "slide_number"
                return "footer"

    # ── 3. Content-based heuristics ─────────────────────
    slide_height = Emu(6858000)  # default 7.5 inches
    try:
        slide_height = slide.slide_layout.slide_master.slide_height
    except (AttributeError, TypeError):
        pass

    shape_top = shape.top if hasattr(shape, "top") and shape.top is not None else Emu(0)

    # 3a. Slide number pattern: pure digit or "N / M" format
    if re.match(r"^\d+\s*/\s*\d+$", text):
        return "slide_number"
    if text.isdigit():
        # Only treat as slide_number if it's short and likely a number
        if len(text) <= 4:
            return "slide_number"

    # 3b. Very large font → section_number (≥ 40pt, usually decorative)
    max_font_size = _get_max_font_size(shape)
    if max_font_size is not None and max_font_size >= 40:
        return "section_number"

    # 3c. Bottom area + small text → footer
    bottom_threshold = slide_height - Emu(914400)  # bottom 1 inch
    if shape_top >= bottom_threshold:
        if max_font_size is not None and max_font_size <= 12:
            return "footer"
        if not text:  # empty text box in footer area
            return "footer"

    # 3d. Top area: distinguish title vs header
    top_threshold = Emu(1143000)  # ~1 inch
    if shape_top < top_threshold:
        # Only classify as title if font is reasonably large (> 14pt)
        # Small text at top is likely a header, not a title
        if max_font_size is not None and max_font_size > 14:
            return "title"
        # Small text at top → likely header/footer
        return "footer"

    # 3e. Small font → caption (≤ 11pt)
    if max_font_size is not None and max_font_size <= 11:
        return "caption"

    # ── 4. Fallback: body ───────────────────────────────
    return "body"


def _get_max_font_size(shape: Any) -> float | None:
    """Get the maximum font size (in pt) across all runs in a shape."""
    if not shape.has_text_frame:
        return None
    max_size = None
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if run.font.size:
                pt = _emu_to_pt(run.font.size)
                if pt is not None:
                    max_size = pt if max_size is None else max(max_size, pt)
    return max_size


def apply_font_fix(run: Any, family: str | None = None, size_pt: float | None = None,
                   bold: bool | None = None, color: str | None = None) -> None:
    """Apply font fixes to a run."""
    font = run.font
    if family:
        font.name = family
    if size_pt is not None:
        font.size = Pt(size_pt)
    if bold is not None:
        font.bold = bold
    if color:
        font.color.rgb = hex_to_rgb(color)


def apply_alignment_fix(para: Any, alignment: str) -> None:
    """Apply alignment fix to a paragraph."""
    if alignment in ALIGNMENT_MAP:
        para.alignment = ALIGNMENT_MAP[alignment]


def apply_line_spacing_fix(para: Any, line_spacing: float) -> None:
    """Apply line spacing fix to a paragraph."""
    para.line_spacing = Pt(line_spacing) if line_spacing else None


def _emu_to_pt(emu: Any) -> float | None:
    """Convert EMU to points."""
    if emu is None:
        return None
    try:
        return float(emu) / 12700
    except (TypeError, ValueError):
        return None


def _align_to_str(align: Any) -> str | None:
    """Convert PP_ALIGN enum to string."""
    if align is None:
        return None
    mapping = {v: k for k, v in ALIGNMENT_MAP.items()}
    return mapping.get(align)


def _align_to_enum(align_str: str | None) -> PP_ALIGN | None:
    """Convert string to PP_ALIGN enum."""
    if align_str is None:
        return None
    return ALIGNMENT_MAP.get(align_str)
