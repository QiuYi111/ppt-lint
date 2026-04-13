"""Lint engine — runs compiled rules against a presentation.

Supports AI-powered role classification via Claude CLI.
Role maps are computed per-slide (batch or one-by-one) and passed
to all checkers, so classification happens once per shape.
"""

from __future__ import annotations

import logging
from pathlib import Path

from pptx import Presentation

from internal.domain.models import CompiledRuleSet, LintIssue, Severity

from .role_classifier import classify_all_slides

logger = logging.getLogger(__name__)


class LintResult:
    """Result of a lint run."""

    def __init__(self, issues: list[LintIssue], file_path: str) -> None:
        self.issues = issues
        self.file_path = file_path

    @property
    def errors(self) -> list[LintIssue]:
        return [i for i in self.issues if i.severity == Severity.ERROR]

    @property
    def warnings(self) -> list[LintIssue]:
        return [i for i in self.issues if i.severity == Severity.WARNING]

    @property
    def infos(self) -> list[LintIssue]:
        return [i for i in self.issues if i.severity == Severity.INFO]

    @property
    def fixable(self) -> list[LintIssue]:
        return [i for i in self.issues if i.fix is not None]

    @property
    def total(self) -> int:
        return len(self.issues)

    @property
    def passed(self) -> bool:
        return self.total == 0


def lint_file(
    file_path: str | Path,
    compiled_rules: CompiledRuleSet,
    use_ai: bool = True,
) -> LintResult:
    """Lint a PPTX file using compiled rules."""
    path = Path(file_path)
    if not path.exists():
        raise FileNotFoundError(f"File not found: {path}")

    prs = Presentation(str(path))
    user_roles = set(compiled_rules.user_roles) if compiled_rules.user_roles else None

    # Classify roles for all slides upfront (batch AI or heuristic)
    all_roles = classify_all_slides(prs, user_roles, use_ai=use_ai)

    all_issues: list[LintIssue] = []
    for slide_index, slide in enumerate(prs.slides):
        role_map = all_roles.get(slide_index, {})

        for checker in compiled_rules.checkers:
            try:
                issues = checker(slide, slide_index, role_map=role_map)
                if issues:
                    all_issues.extend(issues)
            except TypeError:
                try:
                    issues = checker(slide, slide_index)
                    if issues:
                        all_issues.extend(issues)
                except Exception as e:
                    logger.error(f"Error running checker on slide {slide_index}: {e}")
            except Exception as e:
                logger.error(f"Error running checker on slide {slide_index}: {e}")

    return LintResult(issues=all_issues, file_path=str(path))


def fix_file(
    file_path: str | Path,
    compiled_rules: CompiledRuleSet,
    output_path: str | Path | None = None,
    dry_run: bool = False,
    use_ai: bool = True,
) -> LintResult:
    """Lint and fix a PPTX file."""
    path = Path(file_path)
    if not path.exists():
        raise FileNotFoundError(f"File not found: {path}")

    prs = Presentation(str(path))
    user_roles = set(compiled_rules.user_roles) if compiled_rules.user_roles else None
    all_roles = classify_all_slides(prs, user_roles, use_ai=use_ai)

    all_issues: list[LintIssue] = []
    fix_actions: list[tuple[int, object, object]] = []

    for slide_index, slide in enumerate(prs.slides):
        role_map = all_roles.get(slide_index, {})

        for checker in compiled_rules.checkers:
            try:
                issues = checker(slide, slide_index, role_map=role_map)
            except TypeError:
                try:
                    issues = checker(slide, slide_index)
                except Exception as e:
                    logger.error(f"Error on slide {slide_index}: {e}")
                    continue
            except Exception as e:
                logger.error(f"Error on slide {slide_index}: {e}")
                continue

            for issue in issues:
                if issue.fix and not dry_run:
                    fix_actions.append((slide_index, issue, slide))
                all_issues.append(issue)

    applied = 0
    if not dry_run:
        for slide_index, issue, slide in fix_actions:
            try:
                issue.fix.apply()
                applied += 1
            except Exception as e:
                logger.error(f"Failed to apply fix on slide {slide_index}: {e}")

    if not dry_run and applied > 0:
        out = Path(output_path) if output_path else path
        prs.save(str(out))
        logger.info(f"Applied {applied} fixes, saved to {out}")

    return LintResult(issues=all_issues, file_path=str(path))
